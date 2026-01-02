from typing import List

try:
    from pptx import Presentation
except ImportError as e:
    raise ImportError(
        "PPTX support requires python-pptx.\n"
        "Install with: pip install omnidoc-sdk[office]"
    ) from e

from omnidoc.core.document import Document, Section, Table


class PPTXExtractor:
    """
    Production-grade PPTX extractor.

    Guarantees:
    - Slide-aware sections
    - Title-based heading detection
    - Bullet grouping
    - RAG semantic chunk compatibility
    """

    def extract(self, path: str) -> Document:
        prs = Presentation(path)

        sections: List[Section] = []
        tables: List[Table] = []

        page = 0

        for slide in prs.slides:
            page += 1
            slide_lines: List[str] = []
            slide_title: str | None = None

            # -------------------------
            # Shape parsing
            # -------------------------
            for shape in slide.shapes:
                # ---------- Tables ----------
                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        rows.append([cell.text.strip() for cell in row.cells])

                    tables.append(
                        Table(
                            page=page,
                            headers=rows[0] if rows else [],
                            rows=rows[1:] if len(rows) > 1 else [],
                        )
                    )

                # ---------- Text ----------
                if not shape.has_text_frame:
                    continue

                text = shape.text.strip()
                if not text:
                    continue

                # Title placeholder → heading
                if shape == slide.shapes.title:
                    slide_title = text
                else:
                    slide_lines.append(text)

            # -------------------------
            # Build sections
            # -------------------------
            if slide_title:
                sections.append(
                    Section(page=page, text=slide_title)
                )

            if slide_lines:
                sections.append(
                    Section(
                        page=page,
                        text="\n".join(slide_lines)
                    )
                )

        raw_text = "\n\n".join(s.text for s in sections)

        return Document(
            metadata={
                "file": path,
                "pages": page,
                "type": "pptx",
            },
            sections=sections,
            tables=tables,
            raw_text=raw_text,
        )
